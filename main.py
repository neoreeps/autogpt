import os
import json
import fitz
import streamlit as st
from docx import Document
from pptx import Presentation
from streamlit_chat import message
from chatbot import ChatBot
from todoist_repair_agent import parse_base_model_with_retries
from todoist_agent.todoist_action_toolkit import TodoistActionToolKit
from todoist_agent.models import (
    ReactResponse,
    CreateNewProjectAction,
    GetAllInboxTasksAction,
    GetAllProjectsAction,
    GetAllTasksAction,
    GiveFinalAnswerAction,
    MoveTaskAction,
)
from logger import get_logger


def read_text_from_file(file, log):
    """
    Read the contents of a .docx file and return the full text.

    Args:
        file (str): The path to the .docx file.

    Returns:
        str: The full text extracted from the .docx file.
    """
    log.info(f"Reading text from file: {file}")
    file_type = file.name.split(".")[-1].lower()
    if file_type == "docx":
        doc = Document(file)
        full_text = []
        for i in range(0, len(doc.tables)):
            full_text.append(f"\nTable {i+1}")
            table = doc.tables[i]
            for row in table.rows:
                for cell in row.cells:
                    full_text.append(cell.text)
        for para in doc.paragraphs:
            full_text.append(para.text)
        return "\n".join(full_text)
    elif file_type == "pptx":
        prs = Presentation(file)
        full_text = []
        for slide_number, slide in enumerate(prs.slides):
            full_text.append(f"\nSlide {slide_number + 1}")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return "\n".join(full_text)
    elif file_type == "pdf":
        text = ""
        with fitz.open(file, file.read()) as doc:
            for page in doc:
                text += page.get_text()
        return text
    else:
        log.error(f"Unsupported file type: {file_type}")
        return "Unsupported file type"


def todoist_agent_loop(chatbot, user_input, temp, hist_len, max_actions, todoist_api_key):
    """
    Executes a loop of actions for a Todoist agent.

    Args:
        chatbot (Chatbot): The chatbot instance.
        user_input (str): The user's input.
        temp (float): The temperature value for generating responses.
        hist_len (int): The length of chat history to consider.
        max_actions (int): The maximum number of actions to perform.
        todoist_api_key (str): The API key for Todoist.

    Returns:
        None
    """
    chatbot.set_todoist_prompt(ReactResponse, user_input)
    todoist = TodoistActionToolKit(todoist_api_key)

    inputs = json.dumps({"objective": user_input})
    for i in range(max_actions):
        raw_response = chatbot.send('user', inputs, temp, hist_len)
        try:
            response = parse_base_model_with_retries(raw_response, ReactResponse)  # noqa
            message(f"Thought: {response.thought}\n" +
                    f"\nAction: {response.action.dict()}\n" +
                    f"\nNumber of actions used: {i + 1}")

            chatbot.messages.append({"role": "assistant", "content": json.dumps(response.dict())})

            match response.action:
                case GiveFinalAnswerAction():
                    message(f"Final Answer: {response.action.answer}") # noqa
                    return
                case GetAllInboxTasksAction():
                    observation = todoist.get_inbox_tasks()
                case GetAllTasksAction():
                    observation = todoist.get_all_tasks()
                case GetAllProjectsAction():
                    observation = todoist.get_all_projects()
                case MoveTaskAction(task_id=task_id, project_id=project_id):
                    todoist.move_task(task_id, project_id)
                    observation = (f"Task with id {task_id} moved to project with id {project_id}.")
                case CreateNewProjectAction(project_name=project_name):
                    observation = todoist.create_project(project_name)
                case _:
                    raise ValueError(f"Unknown action {response.action}")
        except ValueError as e:
            observation = f"You response caused the following error: {e}. Please try again and avoid this error."
            chatbot.messages.append({"role": "assistant", "content": observation})

        # message(f"Observation: {observation}")
        inputs = json.dumps({"observation": observation})

    message("I have used my maximum number of actions. I will now stop.")


def main() -> None:
    """
    The main function that sets up the Streamlit app and handles user interactions.

    Returns:
        None
    """
    if 'logger' in st.session_state:
        del st.session_state.logger

    # Set up logging
    st.session_state.logger = get_logger("main")
    log = st.session_state.logger

    log.info("Starting main function")

    # Set up the layout of the Streamlit app
    st.set_page_config(page_title="Content GPT Writer", layout="wide")
    st.title("Auto GPT")
    st.write('See the code: https://github.com/neoreeps/autogpt')

    # Add a file uploader for documents
    uploaded_files = []

    # Add a sidebar for settings
    with st.sidebar:
        # Add radio buttons for choosing GPT engine and content type, and a text input for API key
        openai_api_key = os.getenv('OPENAI_API_KEY', None)
        if not openai_api_key:
            st.write("You must provide an OpenAI API key set in the environment.")

        gpt_engine_choice = st.selectbox("Choose GPT engine:", ("gpt-4o", "gpt-4-turbo", "gpt-4-turbo-preview", "gpt-4", "gpt-3.5-turbo"))
        log.debug(f"Selected GPT engine: {gpt_engine_choice}")
        temp = st.slider("Select the temperature (entropy): ", 0.0, 1.0, 0.5)
        hist_len = st.slider("Select the history length:", 1, 50, 25)
        content_type = st.radio("Select the type of agent to use:",
                                ("general", "todoist"))

        welcome = "Ask me anything and I'll do my best." + \
            f"  I remember context up to the last {hist_len} messages." + \
            "  I can generate or improve code and todo lists."

        if content_type == "todoist":
            todoist_api_key = os.getenv('TODOIST_API_KEY', None)
            if not todoist_api_key:
                st.write("You must provide a Todoist API key set in the environment.")
            max_actions = st.slider("Select the maximum number of actions to take:", 1, 300, 50)
            st.write("Todoist integration based largely on https://github.com/j0rd1smit/todoist_react_agent")
            welcome = "Ask me about your todo list or what you'd like to add to it."
        else:
            uploaded_files = st.file_uploader("Upload one or more documents to use in your context",
                                              type=["docx", "pptx", "pdf"], accept_multiple_files=True)
            st.write("Note: The documents will be used in the system prompt labeled as 'DOCUMENT 0', 'DOCUMENT 1', etc.")  # noqa
            welcome = "Ask me anything and I'll do my best."

        if 'chatbot' in st.session_state and gpt_engine_choice != st.session_state.gpt_engine:
            del st.session_state.chatbot

    # Create an instance of the ChatBot class only once
    if 'chatbot' not in st.session_state:
        st.session_state.gpt_engine = gpt_engine_choice
        st.session_state.chatbot = ChatBot(openai_api_key, gpt_engine_choice)

    # Get the instance of the ChatBot class
    chatbot = st.session_state.chatbot

    # Set the system prompt
    ext_prompt = "\nFollow the user's requirements carefully and to the letter."
    for x in range(0, len(uploaded_files)):
        ext_prompt = ext_prompt + f"\nDOCUMENT {x}: {read_text_from_file(uploaded_files[x], log)}"
    chatbot.set_system_prompt(content_type, ext_prompt)

    # Allow the user to update the prompt
    with st.expander("Edit the system prompt below, the default is shown:"):
        prompt = st.text_area("System Prompt:",
                              chatbot.messages[0]["content"],
                              height=200)
        chatbot.set_system_prompt(content_type, prompt)

    # Display the chat history or welcome message
    message(f"Hello Human! {welcome}", is_user=False)
    for msg in chatbot.messages[1:]:
        is_user = True if msg["role"] == "user" else False
        message(msg["content"], is_user)

    user_input = st.chat_input("Type your request here ...")
    if user_input:
        message(user_input, is_user=True)
        with st.spinner("Thinking..."):
            if content_type == "todoist":
                todoist_agent_loop(chatbot, user_input, temp, hist_len, max_actions, todoist_api_key)
                # clear the chat history after each iteration
                chatbot.messages = chatbot.messages[:1]
            else:
                message(chatbot.send("user", user_input, temp, hist_len), is_user=False)
    st.write(f"History Depth: {str(chatbot.messages.__len__())}")

    if st.button("Clear"):
        chatbot.messages = chatbot.messages[:1]
        # Refresh the page to show changes
        st.rerun()


if __name__ == "__main__":
    main()
